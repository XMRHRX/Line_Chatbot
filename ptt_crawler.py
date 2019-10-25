import urllib
import requests
import json
from bs4 import BeautifulSoup as bs
import pandas as pd
import time
import datetime
import re
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
import smtplib   
from email.mime.multipart import MIMEMultipart #email內容載體
from email.mime.text import MIMEText #用於製作文字內文
from email.mime.base import MIMEBase #用於承載附檔
from email import encoders #用於附檔編碼
from email.mime.application import MIMEApplication#傳送附件
import ptt_crawler_index
from ptt_crawler_function import *

'''
下列程式碼的邏輯會分成3個階段：
1.先取得每一篇文章的標題、連結等資訊，
也就是 https://www.ptt.cc/bbs/Gossiping/index.html 這一頁會看到的畫面

2.再進入每一篇文章的網頁，例如
https://www.ptt.cc/bbs/Gossiping/M.1568695082.A.850.html
才能得到文章的內文、詳細發佈時間、推噓文等等資訊。

3.整理所需資料，把資料計算出我們所需要的資訊

4製作ppt

5.寄送Email

'''


####==== 1.先取得每篇文章的標題、連結 ====####
title_text=[]
title_url=[]
date_text=[]
push_text=[]
for i in range(ptt_crawler_index.page_n):
    time.sleep(5)
    url = 'https://www.ptt.cc/bbs/'+ptt_crawler_index.board+'/search?page='+str(i+1)+'&q='+urllib.parse.quote(ptt_crawler_index.keyword)
    r = requests.get(url,cookies={'over18': '1'})
    soup = bs(r.text,"html.parser")

    #處理標題字串
    title_raw= soup.select(ptt_crawler_index.title_node)
    for title in title_raw:
        title_text.append(title.text)
        title_url.append('https://www.ptt.cc'+title['href'])

    #處理日期字串
    date_raw= soup.select(ptt_crawler_index.date_node)

    for date in date_raw:
        date_text.append(date.text)

    #處理推/噓字串
    push_raw= soup.select(ptt_crawler_index.push_node)
    for push in push_raw:
        push_text.append(push.text)


    title_result= pd.DataFrame({
        '標題文字':title_text,
        '標題連結':title_url,
        '日期':date_text,
        '推噓':push_text
    })


####==== 2.進入每一篇文章裡面，取得內文、推噓文 ====####
content_text=[]
push_text=[]
article_time_text=[]
push=[]
arrow=[]
sh=[]

for url in title_result['標題連結']:
    time.sleep(5)
    article_resp = requests.get(url,cookies={'over18': '1'})
    article_soup = bs(article_resp.text,"html.parser")
    content_text.append(article_text(article_soup))
    push_text.append(push_content(article_soup))
    article_time_text.append(article_time(article_soup))
    push_count = push_num(article_soup)
    push.append(push_count[0])
    arrow.append(push_count[1])
    sh.append(push_count[2])
    

total_result = pd.concat([title_result,pd.DataFrame({
    '文章內容':content_text,
    '推文內容':push_text,
    '文章詳細時間':article_time_text,
    '推文數':push,
    '箭頭數':arrow,
    '噓文數':sh
})],axis=1)

total_result.to_excel(ptt_crawler_index.attach_file_name_2)

####==== 3.開始計算文章中的正負向詞彙出現數量(僅計算正文) ====####
pos_word = pd.read_excel(ptt_crawler_index.pos_neg_word_file,sheet_name='正向',header=None)[0]
neg_word = pd.read_excel(ptt_crawler_index.pos_neg_word_file,sheet_name='負向',header=None)[0]
pos_result = count_word(total_result['文章內容'],pos_word)
neg_result = count_word(total_result['文章內容'],neg_word)
pos_sum = pos_result['Total'].sum()
neg_sum = neg_result['Total'].sum()


####==== 4.生成PPT ====####

prs = Presentation() #如果有範本可以放在這裡  Presentation('template.pptx')

	#第一頁：封面：要包含大小標
slide_page_1 = prs.slides.add_slide(prs.slide_layouts[0]) #新增的投影片是哪種版面配置
slide_page_1.shapes.title.text = '輿情分析系統\n-以PTT八卦板為例'   # 設定大標題，\n為換行符號
slide_page_1.placeholders[1].text = '當機器人來上班\n-未來職場的ＡＩ必修課'# 設定小標題，\n為換行符號

	#第二頁：本次執行summary
slide_page_2 = prs.slides.add_slide(prs.slide_layouts[1]) 
slide_page_2.shapes.title.text = '本次執行摘要' #大標題
#在標題下面的文字框新增文字，其實就是一串字串結合起來。
slide_page_2.placeholders[1].text_frame.text =\
'關鍵字： '+ptt_crawler_index.keyword+'\n'+\
'執行時間：'+datetime.datetime.now().strftime('%Y-%m-%d %H:%M')+'\n'+\
'爬取共 ' + str(ptt_crawler_index.page_n) +' 頁，共有 '+str(len(total_result))+' 篇文章，\n'+\
'詞彙字典：共有 '+str(len(pos_word)) +' 個正向詞彙， '+str(len(neg_word))+' 個負向詞彙'


	#第三頁：SHOW 前四筆文章資料，提供使用者對於抓下來的文本有初步概念
slide_page_3 = prs.slides.add_slide(prs.slide_layouts[1])
slide_page_3.shapes.title.text = '最新'+str(ptt_crawler_index.page3_article_show_n)+'筆文章'
table = slide_page_3.shapes.add_table(rows=ptt_crawler_index.page3_article_show_n+1, cols=4, left=Inches(0.1), top=Inches(1.5),width=Inches(5), height=Inches(1.2)).table  #在此修改表格的欄位數，寬度等等
		# 設定單元格寬度
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(2)
table.columns[2].width = Inches(3.5)
table.columns[3].width = Inches(2)

		# 設定標題行
table.cell(0, 0).text = '文章標題'
table.cell(0, 1).text = '發佈時間'
table.cell(0, 2).text = '文章內容'
table.cell(0, 3).text = '連結'

		# 填充資料
for row in range(ptt_crawler_index.page3_article_show_n):
    table.cell(row+1, 0).text = total_result['標題文字'][row]
    table.cell(row+1, 1).text = total_result['文章詳細時間'][row]
    
    if len(total_result['文章內容'][row])>40: #超過40字以下省略
        table.cell(row+1, 2).text = total_result['文章內容'][row][0:40]+'...'
    else:
        table.cell(row+1, 2).text = total_result['文章內容'][row]
        
    table.cell(row+1, 3).text = total_result['標題連結'][row]


	#第四頁：正負向詞彙圓餅圖
if pos_sum>0 or neg_sum>0:
	slide_page_4 = prs.slides.add_slide(prs.slide_layouts[5])
	slide_page_4.shapes.title.text = '正負向詞彙比例圖' #大標題
	chart_data = ChartData()
	chart_data.categories = ['正面','負面']  #所有類別
	chart_data.add_series('詞頻', (pos_sum/(pos_sum+neg_sum),neg_sum/(pos_sum+neg_sum))) #計算後的數值

	x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5) #設定圖片的位置與長寬
	chart = slide_page_4.shapes.add_chart(
		XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
	).chart

	chart.has_legend = True
	chart.legend.position = XL_LEGEND_POSITION.BOTTOM
	chart.legend.include_in_layout = False

	chart.plots[0].has_data_labels = True
	data_labels = chart.plots[0].data_labels
	data_labels.number_format = '0%'
	data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
else:
	slide_page_4 = prs.slides.add_slide(prs.slide_layouts[5])
	slide_page_4.shapes.title.text = '正負向詞彙比例圖 - 資料不足無法繪圖' #大標題


	#第五頁：正負向詞彙次數表  
slide_page_5 = prs.slides.add_slide(prs.slide_layouts[5])
slide_page_5.shapes.title.text = '正負向詞彙次數表(Top'+str(ptt_crawler_index.page5_word_show_n)+')'
table_pos = slide_page_5.shapes.add_table(rows=ptt_crawler_index.page5_word_show_n+1, cols=2, left=Inches(0.5), top=Inches(2),width=Inches(1), height=Inches(1.2)).table
		# 設定單元格寬度
table_pos.columns[0].width = Inches(2)
table_pos.columns[1].width = Inches(2)

    # 設定標題行
table_pos.cell(0, 0).text = '正向詞彙'
table_pos.cell(0, 1).text = '次數'
show_pos_result = pos_result['Total'].head(ptt_crawler_index.page5_word_show_n) 
    # 填充資料
for row in range(ptt_crawler_index.page5_word_show_n):
    table_pos.cell(row+1, 0).text = show_pos_result.index[row]
    table_pos.cell(row+1, 1).text = str(show_pos_result[row])

    
table_neg = slide_page_5.shapes.add_table(rows=ptt_crawler_index.page5_word_show_n+1, cols=2, left=Inches(5), top=Inches(2),width=Inches(1), height=Inches(1.2)).table
    # 設定單元格寬度
table_neg.columns[0].width = Inches(2)# Inches(2.0)
table_neg.columns[1].width = Inches(2)

    # 設定標題行
table_neg.cell(0, 0).text = '負向詞彙'
table_neg.cell(0, 1).text = '次數'
show_neg_result = neg_result['Total'].head(ptt_crawler_index.page5_word_show_n) 
    # 填充資料
for row in range(ptt_crawler_index.page5_word_show_n):
    table_neg.cell(row+1, 0).text = show_neg_result.index[row]
    table_neg.cell(row+1, 1).text = str(show_neg_result[row])


prs.save(ptt_crawler_index.attach_file_name_1)


####==== 5. 寄送email ====####
#使用Gmail要先至https://myaccount.google.com/lesssecureapps 做設定，開啟「允許低安全性應用程式」功能
#gmail的限制為500封信/天 AND 500人/單封

#開始組合信件內容
mail = MIMEMultipart()
mail['From'] = ptt_crawler_index.from_gmail_user
mail['To'] = ptt_crawler_index.to_user
mail['Subject'] = ptt_crawler_index.Subject
#將信件內文加到email中
mail.attach(MIMEText(ptt_crawler_index.contents))     
#將附加檔案們加到email中       
part_attach1 = MIMEApplication(open(ptt_crawler_index.attach_file_name_1,'rb').read())   #開啟附件
part_attach1.add_header('Content-Disposition','attachment',filename=ptt_crawler_index.attach_file_name_1) #為附件命名
mail.attach(part_attach1)   #新增附件

part_attach2 = MIMEApplication(open(ptt_crawler_index.attach_file_name_2,'rb').read())   #開啟附件
part_attach2.add_header('Content-Disposition','attachment',filename=ptt_crawler_index.attach_file_name_2) #為附件命名
mail.attach(part_attach2)


server = smtplib.SMTP_SSL(ptt_crawler_index.which_server,ptt_crawler_index.server_smtp_port)  #如果換信件系統要改伺服器
server.ehlo()
server.login(ptt_crawler_index.from_gmail_user, ptt_crawler_index.from_gmail_password)
server.send_message(mail)
server.quit()

####完成!####